"""
Generate a teleGlobal-styled architecture slide deck (.pptx) for the
FleetTrack / FleetManagement POC.

Mirrors the 5-slide structure of the reference AI Bank Onboarding deck:
  1. Technical Architecture  (3 pillar cards + request-flow band)
  2. As-Built System Architecture  (layered rows)
  3. Technical Highlights & Engineering Rigor  (4 quadrants)
  4. Deployment Architecture (AWS)
  5. Complete Production Architecture - Full Stack  (7 layers)

Run:  python generate_architecture_ppt.py
Output: FleetTrack_Architecture.pptx  (next to this script)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette
NAVY        = RGBColor(0x0C, 0x1E, 0x3C)   # header background
NAVY_DK     = RGBColor(0x08, 0x16, 0x2C)
BLUE        = RGBColor(0x2E, 0x6F, 0xE0)   # accent / titles
BLUE_BAR    = RGBColor(0x2E, 0x6F, 0xE0)
GREEN       = RGBColor(0x1F, 0x9E, 0x63)
GREEN_BAR   = RGBColor(0x2E, 0xA8, 0x66)
PURPLE      = RGBColor(0x6B, 0x46, 0xC1)
ORANGE      = RGBColor(0xE0, 0x8A, 0x1E)
RED         = RGBColor(0xD1, 0x4B, 0x3B)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
INK         = RGBColor(0x1F, 0x2A, 0x3C)   # body text
MUTED       = RGBColor(0x6B, 0x76, 0x88)
CARD_BORDER = RGBColor(0xD9, 0xDF, 0xEA)
CARD_FILL   = RGBColor(0xFF, 0xFF, 0xFF)
BAND_FILL   = RGBColor(0xE9, 0xF0, 0xFC)
ROW_BLUE    = RGBColor(0xEE, 0xF3, 0xFD)
ROW_GREEN   = RGBColor(0xEC, 0xF6, 0xF0)
ROW_PURPLE  = RGBColor(0xF1, 0xEE, 0xFA)
ROW_ORANGE  = RGBColor(0xFD, 0xF3, 0xE6)
ROW_RED     = RGBColor(0xFB, 0xEE, 0xEC)
ROW_GREY    = RGBColor(0xF4, 0xF6, 0xF9)

# 16:9
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SW = prs.slide_width
SH = prs.slide_height


# ---------------------------------------------------------------- helpers
def _no_autofit(tf):
    tf.word_wrap = True


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def set_line(shape, color, w=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(w)


def rounded(slide, x, y, w, h, fill, line=None, line_w=1.0, radius=0.06):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    set_fill(shp, fill)
    if line is None:
        no_line(shp)
    else:
        set_line(shp, line, line_w)
    shp.shadow.inherit = False
    return shp


def rect(slide, x, y, w, h, fill, line=None, line_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    set_fill(shp, fill)
    if line is None:
        no_line(shp)
    else:
        set_line(shp, line, line_w)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    _no_autofit(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.paragraphs[0].alignment = align
    return tb, tf


def add_para(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT,
             first=False, space_after=2, bullet=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    if bullet:
        _set_bullet(p)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return p


def _set_bullet(p):
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
    # clear existing bullet defs
    for tag in ('a:buNone', 'a:buChar', 'a:buAutoNum'):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.set('indent', str(-Emu(Inches(0.14)).emu if False else '-127000'))
    pPr.set('marL', '127000')
    pPr.append(buChar)


def title_block(slide, title, subtitle, accent):
    # header band
    hdr = rect(slide, 0, 0, SW, Inches(1.32), NAVY)
    # accent bar
    rect(slide, 0, Inches(1.32), SW, Inches(0.07), accent)
    # title
    tb, tf = textbox(slide, Inches(0.55), Inches(0.22), Inches(9.05), Inches(0.7))
    add_para(tf, title, 27, WHITE, bold=True, first=True, space_after=0)
    tb2, tf2 = textbox(slide, Inches(0.57), Inches(0.86), Inches(9.05), Inches(0.4))
    add_para(tf2, subtitle, 13, RGBColor(0xBF, 0xCD, 0xE6), first=True, space_after=0)
    _logo(slide)


HERE = os.path.dirname(os.path.abspath(__file__))

# The real teleGlobal logo is used automatically when found. Save the white
# transparent-background PNG as  docs/teleglobal-logo.png  (any of the names
# below also work) and re-run this script.
_LOGO_CANDIDATES = [
    "teleglobal_logo.png", "teleglobal-logo.png", "teleglobal.png",
    "logo.png", "teleglobal_logo.jpg", "teleglobal-logo.jpg",
    "teleglobal-logo.jpeg", "teleglobal-logo.webp",
]

# Logo layout in the navy header: right-aligned, vertically centred.
LOGO_H        = Inches(0.66)    # rendered height
LOGO_RIGHT    = Inches(0.55)    # gap from the slide's right edge
HEADER_H      = Inches(1.32)    # navy band height (keep in sync with title_block)


def _find_logo():
    search_dirs = [
        HERE,                                   # docs/
        os.path.dirname(HERE),                  # project root
        os.path.dirname(os.path.dirname(HERE)),  # workspace root
        os.path.join(os.path.expanduser("~"), "Downloads"),
        os.path.join(os.path.expanduser("~"), "Pictures"),
    ]
    for d in search_dirs:
        for name in _LOGO_CANDIDATES:
            p = os.path.join(d, name)
            if os.path.isfile(p):
                return p
    return None


LOGO_PATH = _find_logo()


def _logo(slide):
    """Place the teleGlobal logo top-right.

    Uses the real logo image when available, otherwise draws a wordmark
    approximation so the deck still renders.
    """
    if LOGO_PATH:
        # Add at a placeholder origin, then right-align and vertically centre
        # using the picture's real aspect ratio (python-pptx scales width from
        # the image when only height is given).
        pic = slide.shapes.add_picture(LOGO_PATH, 0, 0, height=LOGO_H)
        pic.left = int(SW - LOGO_RIGHT - pic.width)
        pic.top = int((HEADER_H - pic.height) / 2)
        return

    # ---- fallback: drawn wordmark ----
    lx = Inches(10.55)
    ly = Inches(0.30)
    # mark circle
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, lx, ly, Inches(0.62), Inches(0.62))
    set_fill(c, NAVY)
    set_line(c, RGBColor(0x3D, 0x7B, 0xE0), 2.0)
    c.shadow.inherit = False
    ci = slide.shapes.add_shape(MSO_SHAPE.OVAL, lx + Inches(0.13), ly + Inches(0.13),
                                Inches(0.36), Inches(0.36))
    set_fill(ci, RGBColor(0x2E, 0x6F, 0xE0))
    no_line(ci)
    ci.shadow.inherit = False
    # wordmark
    tb, tf = textbox(slide, lx + Inches(0.72), ly - Inches(0.02),
                     Inches(2.0), Inches(0.45), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "tele"
    r1.font.size = Pt(22); r1.font.bold = True; r1.font.color.rgb = WHITE
    r1.font.name = "Segoe UI"
    r2 = p.add_run(); r2.text = "Global"
    r2.font.size = Pt(22); r2.font.bold = True
    r2.font.color.rgb = RGBColor(0x4F, 0x8B, 0xE8); r2.font.name = "Segoe UI"
    tb2, tf2 = textbox(slide, lx + Inches(0.72), ly + Inches(0.36),
                       Inches(2.4), Inches(0.24))
    add_para(tf2, "Elevating business to technology", 8,
             RGBColor(0xBF, 0xCD, 0xE6), first=True, space_after=0)


def footer(slide, text):
    tb, tf = textbox(slide, Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.36))
    add_para(tf, text, 9, MUTED, first=True, space_after=0)


def pillar_card(slide, x, y, w, h, heading, heading_color, lines):
    card = rounded(slide, x, y, w, h, CARD_FILL, line=CARD_BORDER, line_w=1.0, radius=0.05)
    tb, tf = textbox(slide, x + Inches(0.18), y + Inches(0.22),
                     w - Inches(0.36), h - Inches(0.4))
    add_para(tf, heading, 15, heading_color, bold=True, align=PP_ALIGN.CENTER,
             first=True, space_after=8)
    for ln in lines:
        add_para(tf, ln, 11, INK, bullet=True, space_after=5)
    return card


def small_box(slide, x, y, w, h, title, sub, title_color=INK, fill=CARD_FILL,
              border=CARD_BORDER, title_size=11, sub_size=8.5):
    rounded(slide, x, y, w, h, fill, line=border, line_w=1.0, radius=0.08)
    tb, tf = textbox(slide, x + Inches(0.06), y, w - Inches(0.12), h,
                     anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf, title, title_size, title_color, bold=True,
             align=PP_ALIGN.CENTER, first=True, space_after=1)
    if sub:
        add_para(tf, sub, sub_size, MUTED, align=PP_ALIGN.CENTER, space_after=0)


def connector(slide, x1, y1, x2, y2, color, w=1.75):
    """Straight connector line between two points (inches as Emu)."""
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color
    cn.line.width = Pt(w)
    return cn


def row_label(slide, x, y, w, h, label, color):
    tb, tf = textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE)
    for i, part in enumerate(label.split("\n")):
        add_para(tf, part, 12, color, bold=True, first=(i == 0), space_after=0)


# ================================================================ SLIDE 1
def slide1():
    s = prs.slides.add_slide(BLANK)
    title_block(s, "FleetTrack \u2014 Technical Architecture",
                "Real-time fleet tracking, crane booking & AI-assisted vendor operations", BLUE_BAR)

    top = Inches(1.78)
    ch = Inches(2.92)
    gap = Inches(0.28)
    cw = (SW - Inches(1.1) - gap * 2) / 3
    x0 = Inches(0.55)

    pillar_card(s, x0, top, cw, ch, "Frontend", BLUE, [
        "Next.js 16 (App Router), TypeScript",
        "Live map \u2014 React Leaflet + MapTiler",
        "Dashboards & charts (Recharts)",
        "AI Chat + Voice console (SSE streaming)",
        "PDF export \u2014 html2canvas + jsPDF",
    ])
    pillar_card(s, x0 + cw + gap, top, cw, ch, "Real-time & AI", GREEN, [
        "WebSocket telemetry push (live GPS)",
        "SSE streaming for LLM chat",
        "Ollama + Qwen 2.5 3B (local, no keys)",
        "Voice agent \u2014 structured transcripts",
        "Alert engine \u2014 overspeed / low-signal",
    ])
    pillar_card(s, x0 + (cw + gap) * 2, top, cw, ch, "Backend & Data", BLUE, [
        "FastAPI gateway, SQLAlchemy async ORM",
        "Pydantic v2 schemas + validation",
        "PostgreSQL 16 (Docker)",
        "Razorpay payments (signature verified)",
        "MQTT broker + GPS simulator feed",
    ])

    band_y = Inches(5.05)
    rounded(s, x0, band_y, SW - Inches(1.1), Inches(1.62), BAND_FILL, radius=0.04)
    tb, tf = textbox(s, x0 + Inches(0.2), band_y + Inches(0.16),
                     SW - Inches(1.5), Inches(1.35))
    add_para(tf, "Request flow", 14, INK, bold=True, align=PP_ALIGN.CENTER,
             first=True, space_after=6)
    add_para(tf, "Operator \u2192 Next.js UI \u2192 FastAPI routes \u2192 service layer "
                 "(equipment \u00b7 bookings \u00b7 payments \u00b7 reports \u00b7 chat/voice) \u2192 "
                 "PostgreSQL; live GPS via MQTT \u2192 WebSocket push to the map",
             11, INK, bullet=True, space_after=4)
    add_para(tf, "AI chat/voice \u2192 fleet context assembled server-side \u2192 Ollama LLM \u2192 "
                 "streamed back over SSE; Razorpay checkout confirmed via verified signature",
             11, INK, bullet=True, space_after=0)

    footer(s, "Phase 1 shipped: tracking, bookings, payments, reports  |  Phase 2: AI chat/voice, "
              "alert engine, geofencing, analytics")


# ================================================================ SLIDE 2
def slide2():
    s = prs.slides.add_slide(BLANK)
    title_block(s, "As-Built System Architecture",
                "The layered flow implemented in the codebase (Next.js + FastAPI)", BLUE_BAR)

    left = Inches(0.55)
    label_w = Inches(1.55)
    content_x = left + label_w + Inches(0.12)
    content_w = SW - content_x - Inches(0.55)
    y = Inches(1.62)

    def band(y, h, fill):
        rounded(s, left, y, SW - Inches(1.1), h, fill, radius=0.03)

    # Clients
    h1 = Inches(0.82)
    band(y, h1, ROW_BLUE)
    row_label(s, left + Inches(0.12), y, label_w, h1, "Clients", BLUE)
    cw = (content_w - Inches(0.2)) / 2
    small_box(s, content_x, y + Inches(0.14), cw, h1 - Inches(0.28),
              "Operator Browser", "Dashboard \u00b7 Map \u00b7 Bookings \u00b7 Chat/Voice")
    small_box(s, content_x + cw + Inches(0.2), y + Inches(0.14), cw, h1 - Inches(0.28),
              "GPS Devices / Simulator", "MQTT telemetry publishers")

    # API Gateway
    y += h1 + Inches(0.12)
    h2 = Inches(0.9)
    band(y, h2, ROW_BLUE)
    row_label(s, left + Inches(0.12), y, label_w, h2, "FastAPI\nGateway", BLUE)
    cw3 = (content_w - Inches(0.4)) / 3
    small_box(s, content_x, y + Inches(0.16), cw3, h2 - Inches(0.32),
              "REST API", "equipment \u00b7 bookings \u00b7 payments \u00b7 reports")
    small_box(s, content_x + cw3 + Inches(0.2), y + Inches(0.16), cw3, h2 - Inches(0.32),
              "Streaming", "WebSocket telemetry \u00b7 SSE chat")
    small_box(s, content_x + (cw3 + Inches(0.2)) * 2, y + Inches(0.16), cw3, h2 - Inches(0.32),
              "AI endpoints", "/chat \u00b7 /voice \u00b7 fleet context")

    # Services
    y += h2 + Inches(0.12)
    h3 = Inches(1.0)
    band(y, h3, ROW_BLUE)
    row_label(s, left + Inches(0.12), y, label_w, h3, "Service\nLayer", BLUE)
    names = ["Equipment", "Bookings", "Payments", "Reports", "Comm / LLM"]
    subs = ["lifecycle FSM", "date-range + FSM", "Razorpay verify", "utilization + PDF", "chat + voice"]
    n = len(names)
    cwn = (content_w - Inches(0.2) * (n - 1)) / n
    for i, (nm, sb) in enumerate(zip(names, subs)):
        small_box(s, content_x + (cwn + Inches(0.2)) * i, y + Inches(0.18),
                  cwn, h3 - Inches(0.36), nm, sb, title_size=10.5, sub_size=8)

    # AI & Realtime
    y += h3 + Inches(0.12)
    h4 = Inches(0.92)
    band(y, h4, ROW_GREEN)
    row_label(s, left + Inches(0.12), y, label_w, h4, "AI &\nReal-time", GREEN)
    a_names = ["Ollama LLM", "Voice Agent", "Alert Engine", "Geofencing", "MQTT Bridge"]
    a_subs = ["Qwen 2.5 3B", "call transcripts", "overspeed/signal", "boundary zones", "GPS \u2192 WS"]
    cwn = (content_w - Inches(0.2) * (len(a_names) - 1)) / len(a_names)
    for i, (nm, sb) in enumerate(zip(a_names, a_subs)):
        small_box(s, content_x + (cwn + Inches(0.2)) * i, y + Inches(0.16),
                  cwn, h4 - Inches(0.32), nm, sb, title_color=GREEN,
                  border=RGBColor(0xC5, 0xE4, 0xD3), title_size=10.5, sub_size=8)

    # Data
    y += h4 + Inches(0.12)
    h5 = Inches(0.9)
    band(y, h5, ROW_PURPLE)
    row_label(s, left + Inches(0.12), y, label_w, h5, "Data &\nExternal", PURPLE)
    d_names = ["PostgreSQL 16", "Razorpay", "MQTT Broker", "GPS Simulator"]
    d_subs = ["fleet + bookings", "payments API", "Mosquitto", "Python feeder"]
    cwn = (content_w - Inches(0.2) * (len(d_names) - 1)) / len(d_names)
    for i, (nm, sb) in enumerate(zip(d_names, d_subs)):
        small_box(s, content_x + (cwn + Inches(0.2)) * i, y + Inches(0.16),
                  cwn, h5 - Inches(0.32), nm, sb, title_color=PURPLE,
                  border=RGBColor(0xD6, 0xCD, 0xEE), title_size=10.5, sub_size=8)

    footer(s, "Flow: browser \u2192 gateway \u2192 service layer \u2192 PostgreSQL; GPS telemetry via MQTT "
              "\u2192 WebSocket; AI chat/voice via Ollama over SSE.")


# ================================================================ SLIDE 3
def slide3():
    s = prs.slides.add_slide(BLANK)
    title_block(s, "Technical Highlights & Engineering Rigor",
                "Design decisions that make the POC production-minded", BLUE_BAR)

    top = Inches(1.72)
    gap = Inches(0.28)
    qw = (SW - Inches(1.1) - gap) / 2
    qh = Inches(2.36)
    x0 = Inches(0.55)

    def quad(x, y, heading, color, lines):
        rounded(s, x, y, qw, qh, CARD_FILL, line=CARD_BORDER, radius=0.04)
        tb, tf = textbox(s, x + Inches(0.22), y + Inches(0.18), qw - Inches(0.44), qh - Inches(0.36))
        add_para(tf, heading, 15, color, bold=True, align=PP_ALIGN.CENTER,
                 first=True, space_after=8)
        for ln in lines:
            add_para(tf, ln, 10.5, INK, bullet=True, space_after=5)

    quad(x0, top, "Reliability & State", BLUE, [
        "Crane lifecycle is a guarded state machine: available \u2192 booked \u2192 working \u2192 repair \u2192 deceased",
        "Payment signature verified server-side before a booking is confirmed",
        "Ollama unreachable \u2192 graceful degraded reply, never a 5xx",
        "50 pytest integration tests cover all endpoints",
    ])
    quad(x0 + qw + gap, top, "Real-time & AI", GREEN, [
        "WebSocket push keeps the map live without polling",
        "SSE streams LLM tokens for responsive chat",
        "LLM runs locally (Qwen 2.5 3B) \u2014 no data leaves the network, no API keys",
        "Fleet context assembled server-side; prompt never trusts raw input blindly",
    ])
    quad(x0, top + qh + gap, "Extensibility", PURPLE, [
        "Service layer is modular \u2014 new domains plug in without touching routing",
        "Voice agent is provider-agnostic (Twilio / Bland.ai ready)",
        "Vendor & product data is data-driven, not hard-coded",
        "MapTiler / Razorpay keys are env-configurable",
    ])
    quad(x0 + qw + gap, top + qh + gap, "Deployment", ORANGE, [
        "Dockerized PostgreSQL 16; backend on uvicorn, frontend on Node",
        "Env-based secrets (Razorpay, Ollama, MapTiler)",
        "Ollama + MQTT broker run as local containers/services",
        "Clear path to ECS/EKS behind an ALB for production",
    ])

    footer(s, "Stack: FastAPI \u00b7 SQLAlchemy \u00b7 Ollama (Qwen 2.5 3B) \u00b7 Next.js 16 \u00b7 PostgreSQL 16 "
              "\u00b7 Razorpay \u00b7 MQTT \u00b7 Leaflet/MapTiler \u00b7 Recharts")


# ================================================================ SLIDE 4
def slide4():
    s = prs.slides.add_slide(BLANK)
    title_block(s, "Deployment Architecture (AWS)",
                "Containerized services on ECS/EKS behind a load balancer", BLUE_BAR)

    # ---- connectors drawn first so the boxes sit on top of the line ends ----
    I = Inches
    # clients -> ALB
    connector(s, I(2.70), I(2.45), I(3.00), I(2.92), BLUE)
    connector(s, I(2.70), I(3.75), I(3.00), I(3.34), BLUE)
    # ALB -> containers
    connector(s, I(4.40), I(2.95), I(4.95), I(2.68), BLUE)
    connector(s, I(4.40), I(3.32), I(4.95), I(3.70), BLUE)
    # backend container -> managed services
    for cy in (2.36, 3.36, 4.36, 5.36):
        connector(s, I(8.45), I(3.70), I(9.00), I(cy), GREEN, w=1.4)

    # left clients
    lx = Inches(0.55)
    small_box(s, lx, Inches(1.95), Inches(2.15), Inches(1.0),
              "Operator Browser", "Next.js SPA", fill=ROW_BLUE, border=CARD_BORDER)
    small_box(s, lx, Inches(3.25), Inches(2.15), Inches(1.0),
              "GPS Devices", "MQTT publishers", fill=ROW_BLUE, border=CARD_BORDER)

    # ALB
    ax = Inches(3.0)
    small_box(s, ax, Inches(2.6), Inches(1.4), Inches(1.05),
              "ALB", "Application\nLoad Balancer", fill=ROW_ORANGE,
              border=RGBColor(0xE8, 0xC9, 0x8F), title_color=ORANGE)

    # ECS cluster
    cx = Inches(4.75)
    cw = Inches(3.9)
    rounded(s, cx, Inches(1.75), cw, Inches(4.05), ROW_BLUE, line=BLUE, line_w=1.25, radius=0.03)
    tb, tf = textbox(s, cx + Inches(0.18), Inches(1.82), cw - Inches(0.3), Inches(0.3))
    add_para(tf, "ECS / EKS Cluster", 12, BLUE, bold=True, first=True, space_after=0)
    small_box(s, cx + Inches(0.2), Inches(2.25), cw - Inches(0.4), Inches(0.85),
              "Frontend Container", "Node 20 \u00b7 Next.js standalone")
    small_box(s, cx + Inches(0.2), Inches(3.25), cw - Inches(0.4), Inches(0.9),
              "Backend Container", "Python 3.12 \u00b7 FastAPI \u00b7 uvicorn")
    small_box(s, cx + Inches(0.2), Inches(4.35), cw - Inches(0.4), Inches(1.2),
              "Agent / LLM Worker", "Ollama (Qwen 2.5 3B) \u00b7 MQTT bridge \u00b7 GPS ingest")

    # right managed services
    rx = Inches(9.0)
    rw = Inches(3.78)
    items = [
        ("Amazon RDS \u2014 PostgreSQL 16", "fleet + bookings data"),
        ("Amazon S3", "report PDFs / assets"),
        ("Secrets Manager", "DB \u00b7 Razorpay \u00b7 API keys"),
        ("ECR \u2014 Container Registry", "backend / frontend images"),
    ]
    ry = Inches(1.95)
    rh = Inches(0.82)
    for nm, sb in items:
        small_box(s, rx, ry, rw, rh, nm, sb, title_color=GREEN,
                  fill=ROW_GREEN, border=RGBColor(0xC5, 0xE4, 0xD3))
        ry += rh + Inches(0.18)

    # external integrations (fills the space under the client boxes)
    ex_y = Inches(4.55)
    rounded(s, Inches(0.55), ex_y, Inches(3.85), Inches(1.25), ROW_ORANGE,
            line=RGBColor(0xE8, 0xC9, 0x8F), line_w=1.0, radius=0.05)
    tb, tf = textbox(s, Inches(0.65), ex_y + Inches(0.12), Inches(3.65), Inches(1.0))
    add_para(tf, "External Integrations", 11.5, ORANGE, bold=True,
             align=PP_ALIGN.CENTER, first=True, space_after=6)
    for t in ("Razorpay \u2014 UPI / cards / netbanking",
              "MapTiler \u2014 map tiles",
              "Twilio / Bland.ai \u2014 voice call-out"):
        add_para(tf, t, 9.5, INK, align=PP_ALIGN.CENTER, space_after=3)

    # observability band
    oy = Inches(6.0)
    rounded(s, Inches(0.55), oy, SW - Inches(1.1), Inches(0.92), ROW_PURPLE,
            line=PURPLE, line_w=1.0, radius=0.04)
    tb, tf = textbox(s, Inches(0.7), oy + Inches(0.06), Inches(3.0), Inches(0.3))
    add_para(tf, "Observability", 11, PURPLE, bold=True, first=True, space_after=0)
    ob = ["CloudWatch \u2014 logs & alarms", "Prometheus / Grafana \u2014 metrics",
          "Health checks \u2014 /health probes"]
    obw = (SW - Inches(1.1) - Inches(3.2)) / 3
    obx = Inches(3.55)
    for i, t in enumerate(ob):
        tb2, tf2 = textbox(s, obx + obw * i, oy + Inches(0.34), obw, Inches(0.4),
                           anchor=MSO_ANCHOR.MIDDLE)
        add_para(tf2, t, 10, INK, align=PP_ALIGN.CENTER, first=True, space_after=0)

    footer(s, "CI builds images \u2192 ECR \u2192 deployed to ECS/EKS. Local LLM keeps fleet data inside "
              "the account; payments via Razorpay.")


# ================================================================ SLIDE 5
def slide5():
    s = prs.slides.add_slide(BLANK)
    title_block(s, "Complete Production Architecture \u2014 Full Stack",
                "End-to-end system: edge, compute, AI, data, integrations, CI/CD", BLUE_BAR)

    left = Inches(0.5)
    label_w = Inches(1.35)
    content_x = left + label_w + Inches(0.1)
    content_w = SW - content_x - Inches(0.5)
    y = Inches(1.55)
    row_gap = Inches(0.08)

    def layer(y, label, lbl_color, band_fill, cells, border, cell_title_color):
        h = Inches(0.72)
        rounded(s, left, y, SW - Inches(1.0), h, band_fill, radius=0.04)
        row_label(s, left + Inches(0.1), y, label_w, h, label, lbl_color)
        n = len(cells)
        g = Inches(0.14)
        cwn = (content_w - g * (n - 1)) / n
        for i, (nm, sb) in enumerate(cells):
            small_box(s, content_x + (cwn + g) * i, y + Inches(0.09), cwn, h - Inches(0.18),
                      nm, sb, title_color=cell_title_color, border=border,
                      title_size=9.5, sub_size=7.5)
        return y + h + row_gap

    y = layer(y, "Edge &\nIngress", ORANGE, ROW_ORANGE, [
        ("CloudFront CDN", "static assets, TLS"),
        ("AWS WAF", "rate limit, OWASP"),
        ("Route 53", "DNS, health checks"),
        ("ALB", "path routing, HTTPS"),
    ], RGBColor(0xE8, 0xC9, 0x8F), ORANGE)

    y = layer(y, "Compute\n(ECS/EKS)", BLUE, ROW_BLUE, [
        ("Frontend", "Next.js 16 \u00b7 Node 20"),
        ("Backend API", "FastAPI \u00b7 Py 3.12"),
        ("Realtime Worker", "MQTT \u2192 WebSocket"),
        ("Ollama Sidecar", "Qwen 2.5 3B local"),
    ], CARD_BORDER, BLUE)

    y = layer(y, "AI & ML", BLUE, ROW_BLUE, [
        ("Ollama LLM", "chat NL understanding"),
        ("Voice Agent", "transcripts + summary"),
        ("Alert Engine", "overspeed / signal"),
        ("Analytics", "utilization scoring"),
    ], CARD_BORDER, BLUE)

    y = layer(y, "Data &\nStorage", GREEN, ROW_GREEN, [
        ("Amazon RDS", "PostgreSQL 16, async"),
        ("Amazon S3", "report PDFs / docs"),
        ("MQTT Broker", "telemetry stream"),
        ("Migrations", "versioned schema"),
    ], RGBColor(0xC5, 0xE4, 0xD3), GREEN)

    y = layer(y, "Security &\nGovernance", RED, ROW_RED, [
        ("Auth / RBAC", "role-based access"),
        ("AWS KMS", "encryption keys"),
        ("Secrets Manager", "DB \u00b7 Razorpay keys"),
        ("Payment Verify", "signature checks"),
    ], RGBColor(0xE8, 0xC3, 0xBE), RED)

    y = layer(y, "Integrations\n& Notify", PURPLE, ROW_PURPLE, [
        ("Razorpay", "UPI / card / netbank"),
        ("Twilio / Bland.ai", "voice call-out"),
        ("MapTiler", "map tiles"),
        ("SNS / SES", "SMS / email alerts"),
    ], RGBColor(0xD6, 0xCD, 0xEE), PURPLE)

    y = layer(y, "Ops &\nCI/CD", INK, ROW_GREY, [
        ("Prometheus + Grafana", "metrics, dashboards"),
        ("CloudWatch", "logs, alarms"),
        ("GitHub Actions", "build \u2192 test \u2192 push"),
        ("Amazon ECR", "image registry"),
    ], CARD_BORDER, INK)

    footer(s, "Full production BOM: 7 layers \u00b7 real-time GPS \u00b7 local LLM (no external AI dependency) "
              "\u00b7 payments via Razorpay \u00b7 all fleet data stays in-account")


slide1()
slide2()
slide3()
slide4()
slide5()

out = os.path.join(HERE, "FleetTrack_Architecture.pptx")
prs.save(out)
print("Logo:", LOGO_PATH if LOGO_PATH else
      "NOT FOUND - drew wordmark fallback. Save the PNG as docs/teleglobal-logo.png "
      "and re-run to embed the real logo.")
print("Saved:", out)
