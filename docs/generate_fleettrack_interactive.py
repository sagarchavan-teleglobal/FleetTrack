"""
Generate a 7-slide INTERACTIVE teleGlobal showcase deck (.pptx) for the
FleetTrack / FleetManagement POC only.

Designed for a big screen (exhibition / stakeholder demo):
  * bold dark theme with vibrant accent bars (reads from across a room)
  * large type, high contrast, generous spacing
  * INTERACTIVE: every slide carries a clickable nav bar; the hub slide has
    two big clickable cards that jump straight into "how it's built" or
    "how it runs in production". All navigation uses real PowerPoint
    slide-jump actions, so it works in Slide Show mode without add-ins.
  * fade transitions between slides

Slides
  1. Hub / title              - two clickable story cards
  2. Technical Architecture   - 3 pillars + request flow
  3. As-Built System          - layered request flow (Clients -> ... -> Data)
  4. Engineering Rigor        - 4 quadrants of design decisions
  5. Deployment (AWS)         - ALB / ECS-EKS / managed services / external
  6. Full Stack               - 7-layer production BOM
  7. Closing                  - impact stats + shared stack + thank you

Run:  python generate_fleettrack_interactive.py
Out:  FleetTrack_Interactive_Showcase.pptx  (next to this script)
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ------------------------------------------------------------------ palette
BG          = RGBColor(0x0A, 0x12, 0x26)   # slide background (deep navy)
BG_PANEL    = RGBColor(0x12, 0x1D, 0x3A)   # card fill
BG_PANEL_2  = RGBColor(0x17, 0x24, 0x46)   # lighter card fill
HEAD_BAR    = RGBColor(0x08, 0x0F, 0x20)   # header strip

WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
INK         = RGBColor(0xEA, 0xF0, 0xFF)   # primary text on dark
MUTED       = RGBColor(0x9F, 0xB0, 0xD0)   # secondary text
FAINT       = RGBColor(0x6B, 0x7C, 0x9C)

BLUE        = RGBColor(0x3D, 0x7B, 0xFF)
CYAN        = RGBColor(0x22, 0xD3, 0xEE)
TEAL        = RGBColor(0x2D, 0xD4, 0xBF)
PURPLE      = RGBColor(0x8B, 0x5C, 0xF6)
PINK        = RGBColor(0xEC, 0x48, 0x99)
AMBER       = RGBColor(0xF5, 0xB3, 0x01)
GREEN       = RGBColor(0x22, 0xC5, 0x8B)

BRD         = RGBColor(0x2A, 0x3A, 0x60)   # card border
BRD_HI      = RGBColor(0x3D, 0x52, 0x84)

# 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

HERE = os.path.dirname(os.path.abspath(__file__))

# nav bar geometry
NAV_Y = Inches(6.86)
NAV_H = Inches(0.42)

# Labels + accent colour for each of the seven slides, used by the nav bar.
NAV_ITEMS = [
    ("Home",         BLUE),
    ("Architecture", CYAN),
    ("As-Built",     TEAL),
    ("Highlights",   PURPLE),
    ("Deployment",   AMBER),
    ("Full Stack",   PINK),
    ("Closing",      GREEN),
]


# ------------------------------------------------------------------ helpers
def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def set_line(shape, color, w=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(w)


def rect(slide, x, y, w, h, fill, line=None, line_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    set_fill(shp, fill)
    no_line(shp) if line is None else set_line(shp, line, line_w)
    shp.shadow.inherit = False
    return shp


def rounded(slide, x, y, w, h, fill, line=None, line_w=1.0, radius=0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    set_fill(shp, fill)
    no_line(shp) if line is None else set_line(shp, line, line_w)
    shp.shadow.inherit = False
    return shp


def oval(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    set_fill(shp, fill)
    no_line(shp)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(3)
    tf.margin_right = Pt(3)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.paragraphs[0].alignment = align
    return tb, tf


def para(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT,
         first=False, space_after=3, bullet=False, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    if bullet:
        _bullet(p)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return p


def _bullet(p):
    pPr = p._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.set("marL", "132000")
    pPr.set("indent", "-132000")
    pPr.append(pPr.makeelement(qn("a:buChar"), {"char": "\u2022"}))


def fill_background(slide, color=BG):
    """Paint the whole slide, then push it behind everything else."""
    bg = rect(slide, 0, 0, SW, SH, color)
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def accent_glow(slide, x, y, w, h, color):
    shp = rounded(slide, x, y, w, h, color, radius=0.5)
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


# ------------------------------------------------------------------ interactivity
# Slide-jump actions can only be assigned once the target slide exists, but
# the nav bar on slide 1 points forward at slides that haven't been built
# yet. Links are queued during the build and wired up by resolve_links().
PENDING_LINKS = []


def link_to(shape, index):
    """Queue `shape` to jump to slide `index` (0-based) when clicked in a show."""
    PENDING_LINKS.append((shape, index))


def resolve_links():
    """Apply every queued slide-jump action. Call after all slides are built."""
    for shape, index in PENDING_LINKS:
        shape.click_action.target_slide = prs.slides[index]
    return len(PENDING_LINKS)


def nav_bar(slide, current):
    """Clickable navigation strip pinned to the bottom of every slide."""
    rect(slide, 0, NAV_Y - Inches(0.06), SW, Inches(0.02), RGBColor(0x1E, 0x2B, 0x4C))

    n = len(NAV_ITEMS)
    gap = Inches(0.10)
    total_w = SW - Inches(1.1)
    bw = (total_w - gap * (n - 1)) / n
    x = Inches(0.55)

    for i, (label, accent) in enumerate(NAV_ITEMS):
        active = (i == current)
        fill = accent if active else RGBColor(0x14, 0x1F, 0x3C)
        pill = rounded(slide, x, NAV_Y, bw, NAV_H, fill,
                       line=None if active else BRD, line_w=1.0, radius=0.42)
        tf = pill.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(2)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(10.5)
        r.font.bold = True
        r.font.name = "Segoe UI"
        r.font.color.rgb = RGBColor(0x06, 0x0C, 0x1A) if active else MUTED
        if not active:
            link_to(pill, i)
        x += bw + gap


def add_transition(slide, kind="fade"):
    """Inject a slide transition (python-pptx has no API for this)."""
    sld = slide._element
    for el in sld.findall(qn("p:transition")):
        sld.remove(el)
    tr = sld.makeelement(qn("p:transition"), {"advClick": "1"})
    tr.append(tr.makeelement(qn("p:" + kind), {}))
    sld.append(tr)


# ------------------------------------------------------------------ building blocks
def slide_head(slide, kicker, kicker_color, title, subtitle=None, accent=BLUE):
    rect(slide, 0, 0, SW, Inches(1.30), HEAD_BAR)
    rect(slide, 0, Inches(1.30), SW, Inches(0.055), accent)
    oval(slide, Inches(0.55), Inches(0.29), Inches(0.15), Inches(0.15), kicker_color)
    tb, tf = textbox(slide, Inches(0.80), Inches(0.20), Inches(8.6), Inches(0.32))
    para(tf, kicker.upper(), 11.5, kicker_color, bold=True, first=True, space_after=0)
    tb, tf = textbox(slide, Inches(0.55), Inches(0.52), Inches(9.3), Inches(0.52))
    para(tf, title, 26, WHITE, bold=True, first=True, space_after=0)
    if subtitle:
        tb, tf = textbox(slide, Inches(0.57), Inches(0.98), Inches(9.3), Inches(0.30))
        para(tf, subtitle, 12.5, MUTED, first=True, space_after=0)


def card(slide, x, y, w, h, accent, heading=None, lines=None, body=None,
         heading_size=15, line_size=11, fill=BG_PANEL, center_heading=False,
         bullets=True):
    rounded(slide, x, y, w, h, fill, line=BRD, line_w=1.0, radius=0.06)
    rect(slide, x + Inches(0.02), y + Inches(0.02), w - Inches(0.04), Inches(0.055), accent)
    tb, tf = textbox(slide, x + Inches(0.18), y + Inches(0.19),
                     w - Inches(0.36), h - Inches(0.34))
    first = True
    if heading:
        para(tf, heading, heading_size, WHITE, bold=True, first=True, space_after=7,
             align=PP_ALIGN.CENTER if center_heading else PP_ALIGN.LEFT)
        first = False
    if body:
        para(tf, body, line_size, MUTED, first=first, space_after=4,
             align=PP_ALIGN.CENTER if center_heading else PP_ALIGN.LEFT)
        first = False
    for ln in (lines or []):
        para(tf, ln, line_size, MUTED, bullet=bullets, first=first, space_after=4)
        first = False
    return tf


def stat_card(slide, x, y, w, h, accent, number, label):
    rounded(slide, x, y, w, h, BG_PANEL, line=BRD, line_w=1.0, radius=0.07)
    rect(slide, x + Inches(0.02), y + Inches(0.02), w - Inches(0.04), Inches(0.07), accent)
    tb, tf = textbox(slide, x + Inches(0.1), y + Inches(0.30), w - Inches(0.2),
                     h - Inches(0.45), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, number, 44, accent, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=6)
    para(tf, label, 11.5, MUTED, align=PP_ALIGN.CENTER, space_after=0)


def cell(slide, x, y, w, h, title, sub, tcolor=WHITE, fill=BG_PANEL_2, border=BRD,
         tsize=10, ssize=7.8):
    rounded(slide, x, y, w, h, fill, line=border, line_w=0.75, radius=0.1)
    tb, tf = textbox(slide, x + Inches(0.05), y, w - Inches(0.10), h,
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, title, tsize, tcolor, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=1)
    if sub:
        para(tf, sub, ssize, MUTED, align=PP_ALIGN.CENTER, space_after=0)


def layer_row(slide, y, h, label, label_color, cells, left=Inches(0.55),
              label_w=Inches(1.42), band=None):
    """One horizontal architecture layer: label + evenly spaced cells."""
    total_w = SW - Inches(1.1)
    rounded(slide, left, y, total_w, h, band or RGBColor(0x0E, 0x18, 0x32),
            line=BRD, line_w=0.75, radius=0.06)
    tb, tf = textbox(slide, left + Inches(0.14), y, label_w, h, anchor=MSO_ANCHOR.MIDDLE)
    for i, part in enumerate(label.split("\n")):
        para(tf, part, 10.5, label_color, bold=True, first=(i == 0), space_after=0)
    cx = left + label_w + Inches(0.16)
    cw_total = total_w - label_w - Inches(0.34)
    n = len(cells)
    gap = Inches(0.10)
    cw = (cw_total - gap * (n - 1)) / n
    for i, (t, s) in enumerate(cells):
        cell(slide, cx + (cw + gap) * i, y + Inches(0.08), cw, h - Inches(0.16),
             t, s, tcolor=label_color)


def footer(slide, text):
    tb, tf = textbox(slide, Inches(0.55), Inches(6.48), Inches(12.2), Inches(0.32))
    para(tf, text, 9, FAINT, first=True, space_after=0)


# ================================================================== SLIDE 1
def slide1_hub():
    s = prs.slides.add_slide(BLANK)
    fill_background(s)

    accent_glow(s, Inches(-1.2), Inches(-1.0), Inches(4.2), Inches(3.0),
                RGBColor(0x16, 0x24, 0x52))
    accent_glow(s, Inches(10.4), Inches(4.6), Inches(4.6), Inches(3.4),
                RGBColor(0x0C, 0x2A, 0x28))

    badge = rounded(s, Inches(0.75), Inches(0.72), Inches(3.35), Inches(0.42),
                    TEAL, radius=0.42)
    tf = badge.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "FleetTrack  \u00b7  teleGlobal"
    r.font.size = Pt(11.5); r.font.bold = True; r.font.color.rgb = RGBColor(0x06, 0x0C, 0x1A)
    r.font.name = "Segoe UI"

    tb, tf = textbox(s, Inches(0.72), Inches(1.36), Inches(11.6), Inches(1.7))
    para(tf, "Real-Time Fleet,", 50, WHITE, bold=True, first=True, space_after=2)
    para(tf, "Built to Scale", 50, CYAN, bold=True, space_after=6)
    tb, tf = textbox(s, Inches(0.75), Inches(3.02), Inches(11.2), Inches(0.46))
    para(tf, "Live GPS tracking, crane booking & AI-assisted vendor operations \u2014 "
             "one production-grade platform. Click a card to jump straight in.",
         14.5, MUTED, first=True, space_after=0)

    # ---- two clickable story cards ----
    cw = Inches(5.75)
    ch = Inches(2.42)
    cy = Inches(3.68)
    x1 = Inches(0.75)
    x2 = Inches(6.83)

    # How it's built -> slide 2 (Architecture)
    b = rounded(s, x1, cy, cw, ch, RGBColor(0x11, 0x22, 0x4A), line=BRD_HI, line_w=1.5,
                radius=0.07)
    rect(s, x1 + Inches(0.02), cy + Inches(0.02), cw - Inches(0.04), Inches(0.08), BLUE)
    link_to(b, 1)
    tb, tf = textbox(s, x1 + Inches(0.34), cy + Inches(0.30), cw - Inches(0.68), Inches(1.9))
    para(tf, "How It's Built", 21, WHITE, bold=True, first=True, space_after=8)
    para(tf, "Next.js 16 frontend, FastAPI backend, WebSocket/SSE real-time layer "
             "and a local LLM chat + voice console \u2014 the full technical design.",
         12, MUTED, space_after=10)
    para(tf, "Architecture  \u00b7  As-Built  \u00b7  Highlights  \u2192", 12, CYAN,
         bold=True, space_after=0)

    # How it runs in production -> slide 5 (Deployment)
    v = rounded(s, x2, cy, cw, ch, RGBColor(0x14, 0x24, 0x1E), line=BRD_HI, line_w=1.5,
                radius=0.07)
    rect(s, x2 + Inches(0.02), cy + Inches(0.02), cw - Inches(0.04), Inches(0.08), GREEN)
    link_to(v, 4)
    tb, tf = textbox(s, x2 + Inches(0.34), cy + Inches(0.30), cw - Inches(0.68), Inches(1.9))
    para(tf, "How It Runs in Production", 21, WHITE, bold=True, first=True, space_after=8)
    para(tf, "Containerized on ECS/EKS behind a load balancer, RDS-backed, "
             "fully observable \u2014 the deployment story on AWS.",
         12, MUTED, space_after=10)
    para(tf, "Deployment  \u00b7  Full Stack  \u2192", 12, GREEN, bold=True, space_after=0)

    footer(s, "Use the bar below to jump between sections at any time \u00b7 "
              "click a card above to start")
    nav_bar(s, 0)
    add_transition(s)


# ================================================================== SLIDE 2
def slide2_architecture():
    s = prs.slides.add_slide(BLANK)
    fill_background(s)
    slide_head(s, "FleetTrack \u00b7 Technical Architecture",
               CYAN,
               "Three Pillars, One Platform",
               "Real-time fleet tracking, crane booking & AI-assisted vendor operations",
               accent=BLUE)

    top = Inches(1.62)
    ch = Inches(2.62)
    gap = Inches(0.22)
    cw = (SW - Inches(1.1) - gap * 2) / 3
    x0 = Inches(0.55)

    card(s, x0, top, cw, ch, BLUE, "Frontend", [
        "Next.js 16 (App Router), TypeScript",
        "Live map \u2014 React Leaflet + MapTiler",
        "Dashboards & charts (Recharts)",
        "AI Chat + Voice console (SSE streaming)",
        "PDF export \u2014 html2canvas + jsPDF",
    ], center_heading=True)

    card(s, x0 + cw + gap, top, cw, ch, CYAN, "Real-time & AI", [
        "WebSocket telemetry push (live GPS)",
        "SSE streaming for LLM chat",
        "Ollama + Qwen 2.5 3B (local, no keys)",
        "Voice agent \u2014 structured transcripts",
        "Alert engine \u2014 overspeed / low-signal",
    ], center_heading=True)

    card(s, x0 + (cw + gap) * 2, top, cw, ch, PURPLE, "Backend & Data", [
        "FastAPI gateway, SQLAlchemy async ORM",
        "Pydantic v2 schemas + validation",
        "PostgreSQL 16 (Docker)",
        "Razorpay payments (signature verified)",
        "MQTT broker + GPS simulator feed",
    ], center_heading=True)

    by = Inches(4.46)
    bh = Inches(1.86)
    rounded(s, x0, by, SW - Inches(1.1), bh, RGBColor(0x0E, 0x18, 0x32),
            line=BRD, line_w=1.0, radius=0.05)
    rect(s, x0 + Inches(0.02), by + Inches(0.02), SW - Inches(1.14), Inches(0.055), TEAL)
    tb, tf = textbox(s, x0 + Inches(0.28), by + Inches(0.20),
                     SW - Inches(1.66), bh - Inches(0.40))
    para(tf, "Request Flow", 14, TEAL, bold=True, align=PP_ALIGN.CENTER,
         first=True, space_after=8)
    para(tf, "Operator \u2192 Next.js UI \u2192 FastAPI routes \u2192 service layer "
             "(equipment \u00b7 bookings \u00b7 payments \u00b7 reports \u00b7 chat/voice) \u2192 "
             "PostgreSQL; live GPS via MQTT \u2192 WebSocket push to the map",
         11.5, MUTED, bullet=True, space_after=6)
    para(tf, "AI chat/voice \u2192 fleet context assembled server-side \u2192 Ollama LLM \u2192 "
             "streamed back over SSE; Razorpay checkout confirmed via verified signature",
         11.5, MUTED, bullet=True, space_after=0)

    footer(s, "Phase 1 shipped: tracking, bookings, payments, reports  |  "
              "Phase 2: AI chat/voice, alert engine, geofencing, analytics")
    nav_bar(s, 1)
    add_transition(s)


# ================================================================== SLIDE 3
def slide3_as_built():
    s = prs.slides.add_slide(BLANK)
    fill_background(s)
    slide_head(s, "FleetTrack \u00b7 As-Built System",
               TEAL,
               "Layered Request Flow",
               "The exact flow implemented in the codebase (Next.js + FastAPI)",
               accent=TEAL)

    y = Inches(1.56)
    h = Inches(0.80)
    step = h + Inches(0.12)

    layer_row(s, y, h, "Clients", BLUE, [
        ("Operator Browser", "Dashboard \u00b7 Map \u00b7 Bookings \u00b7 Chat/Voice"),
        ("GPS Devices / Simulator", "MQTT telemetry publishers"),
    ], band=RGBColor(0x0D, 0x1A, 0x3E))
    y += step

    layer_row(s, y, h, "FastAPI\nGateway", CYAN, [
        ("REST API", "equipment \u00b7 bookings \u00b7 payments \u00b7 reports"),
        ("Streaming", "WebSocket telemetry \u00b7 SSE chat"),
        ("AI endpoints", "/chat \u00b7 /voice \u00b7 fleet context"),
    ], band=RGBColor(0x0A, 0x1F, 0x3A))
    y += step

    layer_row(s, y, h, "Service\nLayer", BLUE, [
        ("Equipment", "lifecycle FSM"),
        ("Bookings", "date-range + FSM"),
        ("Payments", "Razorpay verify"),
        ("Reports", "utilization + PDF"),
        ("Comm / LLM", "chat + voice"),
    ], band=RGBColor(0x0D, 0x1A, 0x3E))
    y += step

    layer_row(s, y, h, "AI &\nReal-time", GREEN, [
        ("Ollama LLM", "Qwen 2.5 3B"),
        ("Voice Agent", "call transcripts"),
        ("Alert Engine", "overspeed/signal"),
        ("Geofencing", "boundary zones"),
        ("MQTT Bridge", "GPS \u2192 WS"),
    ], band=RGBColor(0x0A, 0x22, 0x1A))
    y += step

    layer_row(s, y, h, "Data &\nExternal", PURPLE, [
        ("PostgreSQL 16", "fleet + bookings"),
        ("Razorpay", "payments API"),
        ("MQTT Broker", "Mosquitto"),
        ("GPS Simulator", "Python feeder"),
    ], band=RGBColor(0x1C, 0x12, 0x3C))

    footer(s, "Flow: browser \u2192 gateway \u2192 service layer \u2192 PostgreSQL; GPS "
              "telemetry via MQTT \u2192 WebSocket; AI chat/voice via Ollama over SSE.")
    nav_bar(s, 2)
    add_transition(s)


# ================================================================== SLIDE 4
def slide4_highlights():
    s = prs.slides.add_slide(BLANK)
    fill_background(s)
    slide_head(s, "FleetTrack \u00b7 Engineering Rigor",
               PURPLE,
               "Built Production-Minded",
               "Design decisions that make the POC production-ready",
               accent=PURPLE)

    top = Inches(1.60)
    gap = Inches(0.22)
    qw = (SW - Inches(1.1) - gap) / 2
    qh = Inches(2.36)
    x0 = Inches(0.55)

    card(s, x0, top, qw, qh, BLUE, "Reliability & State", [
        "Crane lifecycle is a guarded state machine: available \u2192 booked \u2192 "
        "working \u2192 repair \u2192 deceased",
        "Payment signature verified server-side before a booking is confirmed",
        "Ollama unreachable \u2192 graceful degraded reply, never a 5xx",
        "50 pytest integration tests cover all endpoints",
    ], center_heading=True)

    card(s, x0 + qw + gap, top, qw, qh, GREEN, "Real-time & AI", [
        "WebSocket push keeps the map live without polling",
        "SSE streams LLM tokens for responsive chat",
        "LLM runs locally (Qwen 2.5 3B) \u2014 no data leaves the network, "
        "no API keys",
        "Fleet context assembled server-side; prompt never trusts raw "
        "input blindly",
    ], center_heading=True)

    card(s, x0, top + qh + gap, qw, qh, PURPLE, "Extensibility", [
        "Service layer is modular \u2014 new domains plug in without "
        "touching routing",
        "Voice agent is provider-agnostic (Twilio / Bland.ai ready)",
        "Vendor & product data is data-driven, not hard-coded",
        "MapTiler / Razorpay keys are env-configurable",
    ], center_heading=True)

    card(s, x0 + qw + gap, top + qh + gap, qw, qh, AMBER, "Deployment", [
        "Dockerized PostgreSQL 16; backend on uvicorn, frontend on Node",
        "Env-based secrets (Razorpay, Ollama, MapTiler)",
        "Ollama + MQTT broker run as local containers/services",
        "Clear path to ECS/EKS behind an ALB for production",
    ], center_heading=True)

    footer(s, "Stack: FastAPI \u00b7 SQLAlchemy \u00b7 Ollama (Qwen 2.5 3B) \u00b7 Next.js 16 "
              "\u00b7 PostgreSQL 16 \u00b7 Razorpay \u00b7 MQTT \u00b7 Leaflet/MapTiler \u00b7 Recharts")
    nav_bar(s, 3)
    add_transition(s)


# ================================================================== SLIDE 5
def slide5_deployment():
    s = prs.slides.add_slide(BLANK)
    fill_background(s)
    slide_head(s, "FleetTrack \u00b7 Deployment",
               AMBER,
               "Cloud Deployment Architecture (AWS)",
               "Containerized services on ECS/EKS behind a load balancer",
               accent=AMBER)

    top = Inches(1.56)
    gap = Inches(0.22)
    ch = Inches(2.85)
    cw = (SW - Inches(1.1) - gap * 2) / 3
    x0 = Inches(0.55)

    card(s, x0, top, cw, ch, BLUE, "Edge & Load Balancing", [
        "Operator Browser \u2014 Next.js SPA",
        "GPS Devices / Simulator \u2014 MQTT publishers",
        "Application Load Balancer \u2014 HTTPS, path routing",
    ], center_heading=True)

    card(s, x0 + cw + gap, top, cw, ch, CYAN, "ECS / EKS Cluster", [
        "Frontend Container \u2014 Node 20, Next.js standalone",
        "Backend Container \u2014 Python 3.12, FastAPI, uvicorn",
        "Agent / LLM Worker \u2014 Ollama (Qwen 2.5 3B), MQTT bridge, "
        "GPS ingest",
    ], center_heading=True)

    card(s, x0 + (cw + gap) * 2, top, cw, ch, GREEN, "Managed AWS Services", [
        "Amazon RDS \u2014 PostgreSQL 16, fleet + bookings",
        "Amazon S3 \u2014 report PDFs / assets",
        "Secrets Manager \u2014 DB \u00b7 Razorpay \u00b7 API keys",
        "Amazon ECR \u2014 container registry",
    ], center_heading=True)

    by = Inches(4.60)
    bh = Inches(1.55)
    bgap = Inches(0.22)
    bw = (SW - Inches(1.1) - bgap) / 2

    card(s, x0, by, bw, bh, AMBER, "External Integrations", [
        "Razorpay \u2014 UPI / cards / netbanking",
        "MapTiler \u2014 map tiles",
        "Twilio / Bland.ai \u2014 voice call-out",
    ], center_heading=True)

    card(s, x0 + bw + bgap, by, bw, bh, PURPLE, "Observability", [
        "CloudWatch \u2014 logs & alarms",
        "Prometheus / Grafana \u2014 metrics",
        "Health checks \u2014 /health probes",
    ], center_heading=True)

    footer(s, "CI builds images \u2192 ECR \u2192 deployed to ECS/EKS. Local LLM keeps "
              "fleet data inside the account; payments via Razorpay.")
    nav_bar(s, 4)
    add_transition(s)


# ================================================================== SLIDE 6
def slide6_full_stack():
    s = prs.slides.add_slide(BLANK)
    fill_background(s)
    slide_head(s, "FleetTrack \u00b7 Full Stack",
               PINK,
               "Complete Production Architecture",
               "Edge, compute, AI, data, integrations, CI/CD \u2014 the full BOM",
               accent=PINK)

    y = Inches(1.56)
    h = Inches(0.60)
    step = h + Inches(0.075)

    layer_row(s, y, h, "Edge &\nIngress", AMBER, [
        ("CloudFront CDN", "static assets, TLS"),
        ("AWS WAF", "rate limit, OWASP"),
        ("Route 53", "DNS, health checks"),
        ("ALB", "path routing, HTTPS"),
    ], band=RGBColor(0x22, 0x1B, 0x0C))
    y += step

    layer_row(s, y, h, "Compute\n(ECS/EKS)", BLUE, [
        ("Frontend", "Next.js 16 \u00b7 Node 20"),
        ("Backend API", "FastAPI \u00b7 Py 3.12"),
        ("Realtime Worker", "MQTT \u2192 WebSocket"),
        ("Ollama Sidecar", "Qwen 2.5 3B local"),
    ], band=RGBColor(0x0D, 0x1A, 0x3E))
    y += step

    layer_row(s, y, h, "AI & ML", CYAN, [
        ("Ollama LLM", "chat NL understanding"),
        ("Voice Agent", "transcripts + summary"),
        ("Alert Engine", "overspeed / signal"),
        ("Analytics", "utilization scoring"),
    ], band=RGBColor(0x0A, 0x1F, 0x3A))
    y += step

    layer_row(s, y, h, "Data &\nStorage", GREEN, [
        ("Amazon RDS", "PostgreSQL 16, async"),
        ("Amazon S3", "report PDFs / docs"),
        ("MQTT Broker", "telemetry stream"),
        ("Migrations", "versioned schema"),
    ], band=RGBColor(0x0A, 0x22, 0x1A))
    y += step

    layer_row(s, y, h, "Security &\nGovernance", PINK, [
        ("Auth / RBAC", "role-based access"),
        ("AWS KMS", "encryption keys"),
        ("Secrets Manager", "DB \u00b7 Razorpay keys"),
        ("Payment Verify", "signature checks"),
    ], band=RGBColor(0x2A, 0x0F, 0x24))
    y += step

    layer_row(s, y, h, "Integrations\n& Notify", PURPLE, [
        ("Razorpay", "UPI / card / netbank"),
        ("Twilio / Bland.ai", "voice call-out"),
        ("MapTiler", "map tiles"),
        ("SNS / SES", "SMS / email alerts"),
    ], band=RGBColor(0x1C, 0x12, 0x3C))
    y += step

    layer_row(s, y, h, "Ops &\nCI/CD", MUTED, [
        ("Prometheus + Grafana", "metrics, dashboards"),
        ("CloudWatch", "logs, alarms"),
        ("GitHub Actions", "build \u2192 test \u2192 push"),
        ("Amazon ECR", "image registry"),
    ], band=RGBColor(0x14, 0x1B, 0x2E))

    footer(s, "Full production BOM: 7 layers \u00b7 real-time GPS \u00b7 local LLM (no "
              "external AI dependency) \u00b7 payments via Razorpay \u00b7 all fleet data "
              "stays in-account")
    nav_bar(s, 5)
    add_transition(s)


# ================================================================== SLIDE 7
def slide7_closing():
    s = prs.slides.add_slide(BLANK)
    fill_background(s)
    slide_head(s, "Closing",
               GREEN,
               "Production-Minded, Not Just a Demo",
               "The impact numbers behind FleetTrack, and what it's built on",
               accent=GREEN)

    sy = Inches(1.58)
    sh_ = Inches(1.72)
    gap = Inches(0.20)
    sw_ = (SW - Inches(1.1) - gap * 3) / 4
    x0 = Inches(0.55)
    stats = [
        (CYAN,  "5",    "Core service domains \u2014 equipment, bookings, payments, "
                          "reports, comms"),
        (GREEN, "50+",  "Pytest integration tests covering every endpoint"),
        (AMBER, "0",    "API keys needed \u2014 the LLM runs locally via Ollama"),
        (PINK,  "Live", "GPS telemetry pushed over WebSocket, no polling"),
    ]
    for i, (c, n, l) in enumerate(stats):
        stat_card(s, x0 + (sw_ + gap) * i, sy, sw_, sh_, c, n, l)

    # ---- shared stack strip ----
    ssy = Inches(3.56)
    ssh = Inches(0.94)
    rounded(s, x0, ssy, SW - Inches(1.1), ssh, RGBColor(0x0E, 0x18, 0x32),
            line=BRD, line_w=1.0, radius=0.07)
    rect(s, x0 + Inches(0.02), ssy + Inches(0.02), SW - Inches(1.14), Inches(0.055), BLUE)
    tb, tf = textbox(s, x0 + Inches(0.26), ssy + Inches(0.13), Inches(1.9),
                     ssh - Inches(0.26), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "Built on", 12, BLUE, bold=True, first=True, space_after=0)
    chips = ["Next.js 16", "FastAPI", "PostgreSQL 16", "Ollama (Qwen 2.5 3B)",
             "Razorpay", "MQTT", "Leaflet / MapTiler", "Recharts"]
    cwid = Inches(1.34)
    cxx = Inches(2.55)
    for i, chp in enumerate(chips):
        cell(s, cxx + (cwid + Inches(0.09)) * i, ssy + Inches(0.26), cwid, Inches(0.42),
             chp, None, tcolor=INK, fill=RGBColor(0x14, 0x26, 0x50), tsize=8.5)

    # ---- closing band ----
    cy = Inches(4.66)
    cbh = Inches(1.26)
    rounded(s, x0, cy, SW - Inches(1.1), cbh, RGBColor(0x0C, 0x2A, 0x22),
            line=BRD_HI, line_w=1.25, radius=0.07)
    tb, tf = textbox(s, x0, cy + Inches(0.16), SW - Inches(1.1), cbh - Inches(0.32),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "Thank You", 30, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True,
         space_after=4)
    para(tf, "Real-time fleet visibility, smarter bookings, one platform  \u00b7  "
             "Questions & Discussion", 13, GREEN, align=PP_ALIGN.CENTER, space_after=0)

    footer(s, "teleGlobal \u2014 Elevating business to technology  |  "
              "click any label in the bar below to revisit a section")
    nav_bar(s, 6)
    add_transition(s)


# ------------------------------------------------------------------ build
BUILDERS = [slide1_hub, slide2_architecture, slide3_as_built,
            slide4_highlights, slide5_deployment, slide6_full_stack,
            slide7_closing]


def build():
    for fn in BUILDERS:
        fn()
    return resolve_links()


if __name__ == "__main__":
    links = build()
    out = os.path.join(HERE, "FleetTrack_Interactive_Showcase.pptx")
    prs.save(out)
    print("Slides:", len(prs.slides._sldIdLst))
    print("Links:", links)
    print("Saved:", out)
