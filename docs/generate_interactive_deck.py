"""
Generate a 6-slide INTERACTIVE teleGlobal showcase deck (.pptx) that combines:
  - AI Bank Onboarding Platform  (multi-agent LLM KYC pipeline)
  - AI Voice Calling Agent       (automated multilingual outbound voice)

Designed for a big screen at an exhibition:
  * bold dark theme with vibrant accent bars (reads from across a hall)
  * large type, high contrast, generous spacing
  * INTERACTIVE: every slide carries a clickable nav bar; the hub slide has
    two big clickable product cards that jump straight to that product's
    slides. All navigation uses real PowerPoint slide-jump actions, so it
    works in Slide Show mode without any add-ins.
  * fade transitions between slides

Slides
  1. Hub / title            - two clickable product cards
  2. Bank - Architecture    - 3 pillars + request flow
  3. Bank - Production      - 7-layer stack + engineering rigor
  4. Voice - Why & What     - impact stats + 4 solution pillars
  5. Voice - How it works   - call flow + decision tiers + languages
  6. Combined value / close - why it matters + stack + thank you

Run:  python generate_interactive_deck.py
Out:  teleGlobal_AI_Showcase_Interactive.pptx  (next to this script)
"""

import copy
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ------------------------------------------------------------------ palette
BG          = RGBColor(0x0A, 0x12, 0x26)   # slide background (deep navy)
BG_PANEL    = RGBColor(0x12, 0x1D, 0x3A)   # card fill
BG_PANEL_2  = RGBColor(0x17, 0x24, 0x46)   # lighter card fill
HEAD_BAR    = RGBColor(0x08, 0x0F, 0x20)   # header strip

WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
INK         = RGBColor(0xEA, 0xF0, 0xFF)   # primary text on dark
MUTED       = RGBColor(0x9F, 0xB0, 0xD0)   # secondary text
FAINT       = RGBColor(0x6B, 0x7C, 0x9C)

BLUE        = RGBColor(0x3D, 0x7B, 0xFF)
CYAN        = RGBColor(0x22, 0xD3, 0xEE)
TEAL        = RGBColor(0x2D, 0xD4, 0xBF)
PURPLE      = RGBColor(0x8B, 0x5C, 0xF6)
PINK        = RGBColor(0xEC, 0x48, 0x99)
AMBER       = RGBColor(0xF5, 0xB3, 0x01)
GREEN       = RGBColor(0x22, 0xC5, 0x8B)

BRD         = RGBColor(0x2A, 0x3A, 0x60)   # card border
BRD_HI      = RGBColor(0x3D, 0x52, 0x84)

# 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

HERE = os.path.dirname(os.path.abspath(__file__))

# nav bar geometry
NAV_Y = Inches(6.86)
NAV_H = Inches(0.42)

# Labels for the six slides used by the nav bar.
NAV_ITEMS = [
    ("Home",        0),
    ("Bank · Arch", 1),
    ("Bank · Prod", 2),
    ("Voice · Why", 3),
    ("Voice · How", 4),
    ("Value",       5),
]


# ------------------------------------------------------------------ helpers
def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def set_line(shape, color, w=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(w)


def rect(slide, x, y, w, h, fill, line=None, line_w=1.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    set_fill(shp, fill)
    no_line(shp) if line is None else set_line(shp, line, line_w)
    shp.shadow.inherit = False
    return shp


def rounded(slide, x, y, w, h, fill, line=None, line_w=1.0, radius=0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    set_fill(shp, fill)
    no_line(shp) if line is None else set_line(shp, line, line_w)
    shp.shadow.inherit = False
    return shp


def oval(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    set_fill(shp, fill)
    no_line(shp)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(3)
    tf.margin_right = Pt(3)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.paragraphs[0].alignment = align
    return tb, tf


def para(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT,
         first=False, space_after=3, bullet=False, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    if bullet:
        _bullet(p)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return p


def _bullet(p):
    pPr = p._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.set("marL", "132000")
    pPr.set("indent", "-132000")
    pPr.append(pPr.makeelement(qn("a:buChar"), {"char": "\u2022"}))


def fill_background(slide, color=BG):
    """Paint the whole slide, then a subtle darker header strip."""
    bg = rect(slide, 0, 0, SW, SH, color)
    # push to back
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def accent_glow(slide, x, y, w, h, color):
    """A soft accent slab used as a decorative glow band."""
    shp = rounded(slide, x, y, w, h, color, radius=0.5)
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


# ------------------------------------------------------------------ logo
_LOGO_CANDIDATES = [
    "teleglobal_logo.png", "teleglobal-logo.png", "teleglobal.png", "logo.png",
]


def find_logo():
    for d in (HERE, os.path.dirname(HERE)):
        for n in _LOGO_CANDIDATES:
            p = os.path.join(d, n)
            if os.path.isfile(p):
                return p
    return None


LOGO = find_logo()
LOGO_H = Inches(0.58)


def put_logo(slide, right=Inches(0.55), top=None, height=LOGO_H):
    if not LOGO:
        tb, tf = textbox(slide, SW - Inches(3.1), Inches(0.28), Inches(2.5), Inches(0.4))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r1 = p.add_run(); r1.text = "tele"
        r1.font.size = Pt(20); r1.font.bold = True; r1.font.color.rgb = WHITE
        r1.font.name = "Segoe UI"
        r2 = p.add_run(); r2.text = "Global"
        r2.font.size = Pt(20); r2.font.bold = True
        r2.font.color.rgb = RGBColor(0x4F, 0x8B, 0xE8); r2.font.name = "Segoe UI"
        return
    pic = slide.shapes.add_picture(LOGO, 0, 0, height=height)
    pic.left = int(SW - right - pic.width)
    pic.top = int(Inches(0.30)) if top is None else int(top)
    return pic


# ------------------------------------------------------------------ interactivity
# Slide-jump actions can only be assigned once the target slide exists, but the
# nav bar on slide 1 points forward at slides that have not been created yet.
# So links are recorded here during the build and wired up afterwards by
# resolve_links().
PENDING_LINKS = []


def link_to(shape, index):
    """Queue `shape` to jump to slide `index` (0-based) when clicked in a show."""
    PENDING_LINKS.append((shape, index))


def resolve_links():
    """Apply every queued slide-jump action. Call after all slides are built."""
    for shape, index in PENDING_LINKS:
        shape.click_action.target_slide = prs.slides[index]
    return len(PENDING_LINKS)


def nav_bar(slide, current):
    """Clickable navigation strip pinned to the bottom of every slide."""
    rect(slide, 0, NAV_Y - Inches(0.06), SW, Inches(0.02), RGBColor(0x1E, 0x2B, 0x4C))

    n = len(NAV_ITEMS)
    gap = Inches(0.10)
    total_w = SW - Inches(1.1)
    bw = (total_w - gap * (n - 1)) / n
    x = Inches(0.55)

    accents = [BLUE, CYAN, CYAN, PURPLE, PURPLE, GREEN]
    for i, (label, target) in enumerate(NAV_ITEMS):
        active = (i == current)
        fill = accents[i] if active else RGBColor(0x14, 0x1F, 0x3C)
        pill = rounded(slide, x, NAV_Y, bw, NAV_H, fill,
                       line=None if active else BRD, line_w=1.0, radius=0.42)
        tf = pill.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(2)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(10.5)
        r.font.bold = True
        r.font.name = "Segoe UI"
        r.font.color.rgb = RGBColor(0x06, 0x0C, 0x1A) if active else MUTED
        if not active:
            link_to(pill, target)
        x += bw + gap


def add_transition(slide, kind="fade", ms="700"):
    """Inject a slide transition (python-pptx has no API for this)."""
    sld = slide._element
    for el in sld.findall(qn("p:transition")):
        sld.remove(el)
    tr = sld.makeelement(qn("p:transition"), {"advClick": "1"})
    tr.append(tr.makeelement(qn("p:" + kind), {}))
    # transition must sit after cSld/clrMapOvr per the schema
    sld.append(tr)


# ------------------------------------------------------------------ building blocks
def slide_head(slide, kicker, kicker_color, title, subtitle=None, accent=BLUE):
    """Header: accent bar + small kicker + big title (+ optional subtitle)."""
    rect(slide, 0, 0, SW, Inches(1.30), HEAD_BAR)
    rect(slide, 0, Inches(1.30), SW, Inches(0.055), accent)
    # accent dot + kicker
    oval(slide, Inches(0.55), Inches(0.29), Inches(0.15), Inches(0.15), kicker_color)
    tb, tf = textbox(slide, Inches(0.80), Inches(0.20), Inches(8.6), Inches(0.32))
    para(tf, kicker.upper(), 11.5, kicker_color, bold=True, first=True, space_after=0)
    # title
    tb, tf = textbox(slide, Inches(0.55), Inches(0.52), Inches(9.3), Inches(0.52))
    para(tf, title, 26, WHITE, bold=True, first=True, space_after=0)
    if subtitle:
        tb, tf = textbox(slide, Inches(0.57), Inches(0.98), Inches(9.3), Inches(0.30))
        para(tf, subtitle, 12.5, MUTED, first=True, space_after=0)
    put_logo(slide)


def card(slide, x, y, w, h, accent, heading=None, lines=None, body=None,
         heading_size=15, line_size=11, fill=BG_PANEL, center_heading=False,
         bullets=True):
    """Standard dark card with a coloured top accent bar."""
    rounded(slide, x, y, w, h, fill, line=BRD, line_w=1.0, radius=0.06)
    rect(slide, x + Inches(0.02), y + Inches(0.02), w - Inches(0.04), Inches(0.055), accent)
    tb, tf = textbox(slide, x + Inches(0.18), y + Inches(0.19),
                     w - Inches(0.36), h - Inches(0.34))
    first = True
    if heading:
        para(tf, heading, heading_size, WHITE, bold=True, first=True, space_after=7,
             align=PP_ALIGN.CENTER if center_heading else PP_ALIGN.LEFT)
        first = False
    if body:
        para(tf, body, line_size, MUTED, first=first, space_after=4,
             align=PP_ALIGN.CENTER if center_heading else PP_ALIGN.LEFT)
        first = False
    for ln in (lines or []):
        para(tf, ln, line_size, MUTED, bullet=bullets, first=first, space_after=4)
        first = False
    return tf


def stat_card(slide, x, y, w, h, accent, number, label):
    rounded(slide, x, y, w, h, BG_PANEL, line=BRD, line_w=1.0, radius=0.07)
    rect(slide, x + Inches(0.02), y + Inches(0.02), w - Inches(0.04), Inches(0.07), accent)
    tb, tf = textbox(slide, x + Inches(0.1), y + Inches(0.30), w - Inches(0.2),
                     h - Inches(0.45), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, number, 46, accent, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=6)
    para(tf, label, 11.5, MUTED, align=PP_ALIGN.CENTER, space_after=0)


def num_step(slide, x, y, w, h, accent, n, title, desc):
    rounded(slide, x, y, w, h, BG_PANEL, line=BRD, line_w=1.0, radius=0.07)
    rect(slide, x + Inches(0.02), y + Inches(0.02), w - Inches(0.04), Inches(0.055), accent)
    d = Inches(0.46)
    oval(slide, x + Inches(0.20), y + Inches(0.26), d, d, accent)
    tb, tf = textbox(slide, x + Inches(0.20), y + Inches(0.26), d, d,
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, str(n), 19, RGBColor(0x06, 0x0C, 0x1A), bold=True,
         align=PP_ALIGN.CENTER, first=True, space_after=0)
    tb, tf = textbox(slide, x + Inches(0.78), y + Inches(0.28), w - Inches(0.95), Inches(0.42))
    para(tf, title, 13, WHITE, bold=True, first=True, space_after=0)
    tb, tf = textbox(slide, x + Inches(0.20), y + Inches(0.85), w - Inches(0.40),
                     h - Inches(1.0))
    para(tf, desc, 10.5, MUTED, first=True, space_after=0)


def cell(slide, x, y, w, h, title, sub, tcolor=WHITE, fill=BG_PANEL_2, border=BRD,
         tsize=10, ssize=7.8):
    rounded(slide, x, y, w, h, fill, line=border, line_w=0.75, radius=0.1)
    tb, tf = textbox(slide, x + Inches(0.05), y, w - Inches(0.10), h,
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, title, tsize, tcolor, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=1)
    if sub:
        para(tf, sub, ssize, MUTED, align=PP_ALIGN.CENTER, space_after=0)


def layer_row(slide, y, h, label, label_color, cells, left=Inches(0.55),
              label_w=Inches(1.42), band=None):
    """One horizontal architecture layer: label + evenly spaced cells."""
    total_w = SW - Inches(1.1)
    rounded(slide, left, y, total_w, h, band or RGBColor(0x0E, 0x18, 0x32),
            line=BRD, line_w=0.75, radius=0.06)
    tb, tf = textbox(slide, left + Inches(0.14), y, label_w, h, anchor=MSO_ANCHOR.MIDDLE)
    for i, part in enumerate(label.split("\n")):
        para(tf, part, 10.5, label_color, bold=True, first=(i == 0), space_after=0)
    cx = left + label_w + Inches(0.16)
    cw_total = total_w - label_w - Inches(0.34)
    n = len(cells)
    gap = Inches(0.10)
    cw = (cw_total - gap * (n - 1)) / n
    for i, (t, s) in enumerate(cells):
        cell(slide, cx + (cw + gap) * i, y + Inches(0.08), cw, h - Inches(0.16),
             t, s, tcolor=label_color)


def footer(slide, text):
    tb, tf = textbox(slide, Inches(0.55), Inches(6.48), Inches(12.2), Inches(0.32))
    para(tf, text, 9, FAINT, first=True, space_after=0)


def arrow(slide, x, y, w, h, color, rot=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    set_fill(shp, color)
    no_line(shp)
    shp.shadow.inherit = False
    if rot:
        shp.rotation = rot
    return shp


# ================================================================== SLIDE 1
def slide1_hub():
    s = prs.slides.add_slide(BLANK)
    fill_background(s)

    # decorative glows
    accent_glow(s, Inches(-1.2), Inches(-1.0), Inches(4.2), Inches(3.0),
                RGBColor(0x16, 0x24, 0x52))
    accent_glow(s, Inches(10.4), Inches(4.6), Inches(4.6), Inches(3.4),
                RGBColor(0x1C, 0x1B, 0x4A))

    put_logo(s, top=Inches(0.42), height=Inches(0.66))

    # eyebrow badge
    badge = rounded(s, Inches(0.75), Inches(0.72), Inches(2.85), Inches(0.42),
                    PINK, radius=0.42)
    tf = badge.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "teleGlobal  ·  AI PLATFORMS"
    r.font.size = Pt(11.5); r.font.bold = True; r.font.color.rgb = WHITE
    r.font.name = "Segoe UI"

    tb, tf = textbox(s, Inches(0.72), Inches(1.36), Inches(11.6), Inches(1.7))
    para(tf, "Intelligent Automation,", 50, WHITE, bold=True, first=True, space_after=2)
    para(tf, "Built End-to-End", 50, CYAN, bold=True, space_after=6)
    tb, tf = textbox(s, Inches(0.75), Inches(3.02), Inches(11.2), Inches(0.46))
    para(tf, "Two production-grade AI platforms, one engineering standard. "
             "Click a card to jump straight into that story.",
         14.5, MUTED, first=True, space_after=0)

    # ---- two clickable product cards ----
    cw = Inches(5.75)
    ch = Inches(2.42)
    cy = Inches(3.68)
    x1 = Inches(0.75)
    x2 = Inches(6.83)

    # Bank card -> slide 2 (index 1)
    b = rounded(s, x1, cy, cw, ch, RGBColor(0x11, 0x22, 0x4A), line=BRD_HI, line_w=1.5,
                radius=0.07)
    rect(s, x1 + Inches(0.02), cy + Inches(0.02), cw - Inches(0.04), Inches(0.08), BLUE)
    link_to(b, 1)
    tb, tf = textbox(s, x1 + Inches(0.34), cy + Inches(0.30), cw - Inches(0.68), Inches(1.9))
    para(tf, "AI Bank Onboarding Platform", 21, WHITE, bold=True, first=True, space_after=8)
    para(tf, "Multi-agent, LLM-orchestrated self-service KYC pipeline. "
             "OCR, KYC, fraud, face match and product recommendations, "
             "coordinated by a LangGraph supervisor.", 12, MUTED, space_after=10)
    para(tf, "Architecture  ·  Production stack  \u2192", 12, CYAN, bold=True, space_after=0)

    # Voice card -> slide 4 (index 3)
    v = rounded(s, x2, cy, cw, ch, RGBColor(0x24, 0x14, 0x46), line=BRD_HI, line_w=1.5,
                radius=0.07)
    rect(s, x2 + Inches(0.02), cy + Inches(0.02), cw - Inches(0.04), Inches(0.08), PURPLE)
    link_to(v, 3)
    tb, tf = textbox(s, x2 + Inches(0.34), cy + Inches(0.30), cw - Inches(0.68), Inches(1.9))
    para(tf, "AI Voice Calling Agent", 21, WHITE, bold=True, first=True, space_after=8)
    para(tf, "Human-like phone conversations, fully automated. Multilingual, "
             "real-time, cloud-native and infinitely scalable outbound calling.",
         12, MUTED, space_after=10)
    para(tf, "Why it matters  ·  How it works  \u2192", 12, PINK, bold=True, space_after=0)

    footer(s, "Use the bar below to jump between sections at any time \u00b7 "
              "click a card above to start")
    nav_bar(s, 0)
    add_transition(s)


# ================================================================== SLIDE 2
def slide2_bank_arch():
    s = prs.slides.add_slide(BLANK)
    fill_background(s)
    slide_head(s, "AI Bank Onboarding  \u00b7  Platform 1 of 2",
               CYAN,
               "Technical Architecture",
               "Multi-agent, LLM-orchestrated self-service KYC pipeline",
               accent=BLUE)

    top = Inches(1.62)
    ch = Inches(2.62)
    gap = Inches(0.22)
    cw = (SW - Inches(1.1) - gap * 2) / 3
    x0 = Inches(0.55)

    card(s, x0, top, cw, ch, BLUE, "Frontend", [
        "Next.js 14 (App Router), React",
        "ChatWindow (Aria) \u2014 conversational UI",
        "Recommendation Panel \u2014 plain-language results",
        "Admin Dashboard \u2014 JWT-protected, RBAC-gated",
    ], center_heading=True)

    card(s, x0 + cw + gap, top, cw, ch, CYAN, "Agent Layer (LangGraph)", [
        "Supervisor Agent \u2014 stateful routing graph",
        "OCR \u00b7 KYC \u00b7 Fraud \u00b7 Face \u00b7 Recommendation",
        "Conversational Agent (Aria) via Ollama LLM",
        "BaseAgent contract: run() + can_handle()",
    ], center_heading=True)

    card(s, x0 + (cw + gap) * 2, top, cw, ch, PURPLE, "Backend & Data", [
        "FastAPI gateway, SQLAlchemy async ORM",
        "PostgreSQL (Alembic), SQLite dev fallback",
        "AWS S3 (KMS-encrypted), Textract / Rekognition",
        "Prometheus + Grafana, MLflow registry",
    ], center_heading=True)

    # request-flow band
    by = Inches(4.46)
    bh = Inches(1.86)
    rounded(s, x0, by, SW - Inches(1.1), bh, RGBColor(0x0E, 0x18, 0x32),
            line=BRD, line_w=1.0, radius=0.05)
    rect(s, x0 + Inches(0.02), by + Inches(0.02), SW - Inches(1.14), Inches(0.055), TEAL)
    tb, tf = textbox(s, x0 + Inches(0.28), by + Inches(0.20),
                     SW - Inches(1.66), bh - Inches(0.40))
    para(tf, "Request Flow", 14, TEAL, bold=True, align=PP_ALIGN.CENTER,
         first=True, space_after=8)
    para(tf, "Customer (no login) \u2192 FastAPI public routes \u2192 Supervisor Agent \u2192 "
             "specialist agent(s) \u2192 shared FullOnboardingState persisted to "
             "PostgreSQL after every step", 11.5, MUTED, bullet=True, space_after=6)
    para(tf, "Bank staff (JWT + RBAC) \u2192 protected routes \u2192 the same agent layer, "
             "for review and manual re-run", 11.5, MUTED, bullet=True, space_after=6)
    para(tf, "Aria drives the conversation; agent_trace is append-only, giving a "
             "full audit trail of every decision", 11.5, MUTED, bullet=True, space_after=0)

    footer(s, "Phase 1\u20132 shipped: OCR, KYC, fraud, face match  |  Phase 3\u20136: "
              "recommendations, conversational agent, multi-agent orchestration, "
              "production hardening")
    nav_bar(s, 1)
    add_transition(s)


# ================================================================== SLIDE 3
def slide3_bank_prod():
    s = prs.slides.add_slide(BLANK)
    fill_background(s)
    slide_head(s, "AI Bank Onboarding  \u00b7  Production",
               CYAN,
               "Complete Production Architecture",
               "Seven layers, 30+ managed services, zero external LLM dependency",
               accent=CYAN)

    y = Inches(1.56)
    h = Inches(0.60)
    step = h + Inches(0.075)

    layer_row(s, y, h, "Edge &\nIngress", AMBER, [
        ("CloudFront CDN", "static assets, TLS"),
        ("AWS WAF", "rate limit, OWASP"),
        ("Route 53", "DNS, failover"),
        ("API Gateway / ALB", "path routing, HTTPS"),
    ], band=RGBColor(0x22, 0x1B, 0x0C))
    y += step

    layer_row(s, y, h, "Compute\n(ECS/EKS)", BLUE, [
        ("Frontend", "Next.js 14 \u00b7 Node 20"),
        ("Backend API", "FastAPI \u00b7 Py 3.12"),
        ("Agent Worker", "LangGraph supervisor"),
        ("Ollama Sidecar", "Qwen2.5:7b, local"),
    ], band=RGBColor(0x0D, 0x1A, 0x3E))
    y += step

    layer_row(s, y, h, "AI & ML", CYAN, [
        ("LangGraph", "agent orchestration"),
        ("Ollama LLM", "NL understanding"),
        ("OpenCV DNN", "YuNet + SFace"),
        ("scikit-learn", "recommendation ML v1"),
        ("MLflow", "model registry"),
    ], band=RGBColor(0x0A, 0x1F, 0x3A))
    y += step

    layer_row(s, y, h, "Data &\nStorage", GREEN, [
        ("Amazon RDS", "PostgreSQL 16, async"),
        ("Amazon S3", "docs SSE-KMS, 1h URLs"),
        ("AgentSession Store", "state + checkpoint"),
        ("Alembic", "versioned schema"),
    ], band=RGBColor(0x0A, 0x22, 0x1A))
    y += step

    layer_row(s, y, h, "Security &\nGovernance", PINK, [
        ("JWT + RBAC", "4 roles"),
        ("AWS KMS", "S3 \u00b7 RDS \u00b7 Secrets"),
        ("Secrets Manager", "creds \u00b7 JWT \u00b7 tokens"),
        ("CloudTrail", "API audit log"),
        ("IAM", "least-privilege"),
    ], band=RGBColor(0x2A, 0x0F, 0x24))
    y += step

    layer_row(s, y, h, "Integrations\n& Notify", PURPLE, [
        ("AWS Textract", "document OCR"),
        ("AWS Rekognition", "face compare"),
        ("KYC Partners", "Surepass \u00b7 IDfy \u00b7 DigiLocker"),
        ("SNS / SES", "SMS OTP \u00b7 email"),
        ("AML / Sanctions", "PEP screening"),
    ], band=RGBColor(0x1C, 0x12, 0x3C))
    y += step

    layer_row(s, y, h, "Ops &\nCI/CD", MUTED, [
        ("Prometheus + Grafana", "metrics, alerting"),
        ("CloudWatch", "logs, anomalies"),
        ("GitHub Actions", "build \u2192 test \u2192 push"),
        ("Amazon ECR", "container registry"),
    ], band=RGBColor(0x14, 0x1B, 0x2E))

    footer(s, "Engineering rigor: forward-only state machine \u00b7 append-only agent_trace "
              "\u00b7 1 retry then graceful degrade \u00b7 20-hop loop guard \u00b7 local LLM keeps "
              "customer data inside the bank's AWS account")
    nav_bar(s, 2)
    add_transition(s)


# ================================================================== SLIDE 4
def slide4_voice_why():
    s = prs.slides.add_slide(BLANK)
    fill_background(s)
    slide_head(s, "AI Voice Calling Agent  \u00b7  Platform 2 of 2",
               PINK,
               "One Agent. Every Conversation.",
               "Human-like phone conversations, fully automated \u2014 the measurable impact",
               accent=PURPLE)

    # ---- impact stats ----
    sy = Inches(1.58)
    sh_ = Inches(1.72)
    gap = Inches(0.20)
    sw_ = (SW - Inches(1.1) - gap * 3) / 4
    x0 = Inches(0.55)
    stats = [
        (CYAN,  "24/7",   "Always-on calling capacity"),
        (TEAL,  "1000s",  "Calls handled in parallel"),
        (AMBER, "~90%",   "Lower cost vs human teams"),
        (PINK,  "100%",   "Consistent script & quality"),
    ]
    for i, (c, n, l) in enumerate(stats):
        stat_card(s, x0 + (sw_ + gap) * i, sy, sw_, sh_, c, n, l)

    # ---- the four solution pillars ----
    py = Inches(3.52)
    ph = Inches(1.60)
    pillars = [
        (BLUE,   1, "Autonomous Calling", "Dials out on its own via cloud telephony."),
        (CYAN,   2, "Natural Voice",      "Understands speech, replies human-like."),
        (PURPLE, 3, "Multilingual",       "Switches languages mid-call, seamlessly."),
        (PINK,   4, "Structured Data",    "Turns every answer into clean data."),
    ]
    for i, (c, n, t, d) in enumerate(pillars):
        num_step(s, x0 + (sw_ + gap) * i, py, sw_, ph, c, n, t, d)

    # ---- problem strip ----
    qy = Inches(5.32)
    qh = Inches(1.00)
    rounded(s, x0, qy, SW - Inches(1.1), qh, RGBColor(0x2A, 0x0F, 0x24),
            line=BRD, line_w=1.0, radius=0.07)
    rect(s, x0 + Inches(0.02), qy + Inches(0.02), SW - Inches(1.14), Inches(0.055), PINK)
    tb, tf = textbox(s, x0 + Inches(0.26), qy + Inches(0.16),
                     SW - Inches(1.62), qh - Inches(0.32))
    para(tf, "Why manual calling doesn't scale", 12.5, PINK, bold=True, first=True,
         space_after=5)
    para(tf, "Limited capacity  \u00b7  expensive teams to hire and retain  \u00b7  "
             "inconsistent tone and data capture  \u00b7  language gaps needing "
             "specialist staff  \u00b7  call insights rarely captured as clean data",
         11.5, MUTED, space_after=0)

    footer(s, "24/7 capacity \u00b7 consistent quality \u00b7 lower cost \u00b7 richer data "
              "\u2014 with no extra headcount")
    nav_bar(s, 3)
    add_transition(s)


# ================================================================== SLIDE 5
def slide5_voice_how():
    s = prs.slides.add_slide(BLANK)
    fill_background(s)
    slide_head(s, "AI Voice Calling Agent  \u00b7  How it works",
               PINK,
               "Architecture, Decisioning & Languages",
               "From API trigger to structured data \u2014 and how the agent picks every reply",
               accent=PURPLE)

    # ---------- architecture flow (left) ----------
    lx = Inches(0.55)
    lw = Inches(7.30)
    rounded(s, lx, Inches(1.56), lw, Inches(2.62), RGBColor(0x0E, 0x18, 0x32),
            line=BRD, line_w=1.0, radius=0.05)
    rect(s, lx + Inches(0.02), Inches(1.58), lw - Inches(0.04), Inches(0.055), CYAN)
    tb, tf = textbox(s, lx + Inches(0.24), Inches(1.72), lw - Inches(0.5), Inches(0.3))
    para(tf, "Solution Architecture", 12.5, CYAN, bold=True, first=True, space_after=0)

    # top row of nodes
    ny = Inches(2.12)
    nh = Inches(0.62)
    nw = Inches(1.94)
    g = Inches(0.28)
    nodes = [("Trigger / UI", "Start Call API"),
             ("Cloud Telephony", "Twilio / Connect"),
             ("Customer", "Phone")]
    for i, (t, sb) in enumerate(nodes):
        nx = lx + Inches(0.28) + (nw + g) * i
        cell(s, nx, ny, nw, nh, t, sb, tcolor=WHITE, fill=RGBColor(0x14, 0x26, 0x50),
             tsize=11, ssize=8.5)
        if i < 2:
            arrow(s, nx + nw + Inches(0.045), ny + Inches(0.20), Inches(0.19),
                  Inches(0.22), CYAN)

    # orchestrator
    oy = Inches(3.00)
    orw = lw - Inches(0.56)
    cell(s, lx + Inches(0.28), oy, orw, Inches(0.56),
         "Application Server \u2014 Orchestrator",
         "Call Flow  \u00b7  State Machine  \u00b7  Language Detection",
         tcolor=WHITE, fill=RGBColor(0x24, 0x18, 0x4E), tsize=11.5, ssize=8.5)
    arrow(s, lx + Inches(3.40), Inches(2.80), Inches(0.19), Inches(0.18), CYAN, rot=90)

    # leaf services
    ly2 = Inches(3.70)
    leaves = [("Speech-to-Text", CYAN), ("Text-to-Speech", AMBER),
              ("LLM (AI Brain)", PURPLE), ("Database", GREEN)]
    lwid = (orw - Inches(0.30)) / 4
    for i, (t, c) in enumerate(leaves):
        px = lx + Inches(0.28) + (lwid + Inches(0.10)) * i
        cell(s, px, ly2, lwid, Inches(0.42), t, None, tcolor=c,
             fill=RGBColor(0x10, 0x1B, 0x38), tsize=9.5)

    # ---------- decision tiers (right) ----------
    rx = Inches(8.10)
    rw = SW - rx - Inches(0.55)
    rounded(s, rx, Inches(1.56), rw, Inches(2.62), RGBColor(0x0E, 0x18, 0x32),
            line=BRD, line_w=1.0, radius=0.05)
    rect(s, rx + Inches(0.02), Inches(1.58), rw - Inches(0.04), Inches(0.055), AMBER)
    tb, tf = textbox(s, rx + Inches(0.24), Inches(1.72), rw - Inches(0.5), Inches(0.3))
    para(tf, "How it decides what to say", 12.5, AMBER, bold=True, first=True, space_after=0)

    tiers = [(GREEN, "1", "Scripted", "Pre-cached questions, ~40 ms", "FASTEST"),
             (AMBER, "2", "FAQ Match", "Keyword answers: price, location", "FAST"),
             (PURPLE, "3", "AI (LLM)", "Dynamic reply for anything off-script", "FLEXIBLE")]
    ty = Inches(2.12)
    th = Inches(0.60)
    for c, n, t, d, tag in tiers:
        rounded(s, rx + Inches(0.24), ty, rw - Inches(0.48), th,
                RGBColor(0x14, 0x1F, 0x3C), line=BRD, line_w=0.75, radius=0.12)
        rect(s, rx + Inches(0.24), ty + Inches(0.06), Inches(0.075), th - Inches(0.12), c)
        d2 = Inches(0.34)
        oval(s, rx + Inches(0.44), ty + Inches(0.13), d2, d2, c)
        tb, tf = textbox(s, rx + Inches(0.44), ty + Inches(0.13), d2, d2,
                         anchor=MSO_ANCHOR.MIDDLE)
        para(tf, n, 13, RGBColor(0x06, 0x0C, 0x1A), bold=True, align=PP_ALIGN.CENTER,
             first=True, space_after=0)
        tb, tf = textbox(s, rx + Inches(0.90), ty + Inches(0.07), rw - Inches(2.30),
                         th - Inches(0.14))
        para(tf, t, 11.5, WHITE, bold=True, first=True, space_after=1)
        para(tf, d, 9, MUTED, space_after=0)
        pill = rounded(s, rx + rw - Inches(1.40), ty + Inches(0.15), Inches(1.10),
                       Inches(0.30), c, radius=0.42)
        tfp = pill.text_frame
        tfp.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = tfp.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        rr = pp.add_run(); rr.text = tag
        rr.font.size = Pt(8); rr.font.bold = True
        rr.font.color.rgb = RGBColor(0x06, 0x0C, 0x1A); rr.font.name = "Segoe UI"
        ty += th + Inches(0.12)

    # ---------- call lifecycle (bottom strip) ----------
    by = Inches(4.42)
    bh = Inches(0.86)
    steps = ["Initiate", "Connect", "Greet", "Listen", "Decide", "Respond", "Loop", "Close"]
    subs = ["API trigger", "customer answers", "first question", "speech to text",
            "pick reply", "text to voice", "each turn", "save & end"]
    n = len(steps)
    gap2 = Inches(0.14)
    bw = (SW - Inches(1.1) - gap2 * (n - 1)) / n
    for i in range(n):
        bx = Inches(0.55) + (bw + gap2) * i
        c = [BLUE, BLUE, CYAN, CYAN, TEAL, AMBER, PURPLE, PINK][i]
        rounded(s, bx, by, bw, bh, RGBColor(0x12, 0x1D, 0x3A), line=BRD, line_w=0.75,
                radius=0.1)
        rect(s, bx + Inches(0.02), by + Inches(0.02), bw - Inches(0.04), Inches(0.05), c)
        tb, tf = textbox(s, bx + Inches(0.05), by + Inches(0.12), bw - Inches(0.10),
                         bh - Inches(0.20), anchor=MSO_ANCHOR.MIDDLE)
        para(tf, f"{i+1}.  {steps[i]}", 10.5, WHITE, bold=True, align=PP_ALIGN.CENTER,
             first=True, space_after=2)
        para(tf, subs[i], 8, MUTED, align=PP_ALIGN.CENTER, space_after=0)

    # ---------- languages ----------
    gy = Inches(5.46)
    gh = Inches(0.86)
    rounded(s, Inches(0.55), gy, SW - Inches(1.1), gh, RGBColor(0x0A, 0x1F, 0x3A),
            line=BRD, line_w=1.0, radius=0.08)
    rect(s, Inches(0.57), gy + Inches(0.02), SW - Inches(1.14), Inches(0.055), TEAL)
    tb, tf = textbox(s, Inches(0.80), gy + Inches(0.14), Inches(3.1), gh - Inches(0.28),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "Speaks their language", 12, TEAL, bold=True, first=True, space_after=2)
    para(tf, "switches mid-call, no interruption", 9, MUTED, space_after=0)
    langs = ["English", "\u0939\u093f\u0928\u094d\u0926\u0940 Hindi",
             "\u092e\u0930\u093e\u0920\u0940 Marathi", "+ More"]
    lwid2 = Inches(1.35)
    lxx = Inches(4.10)
    for i, lg in enumerate(langs):
        cell(s, lxx + (lwid2 + Inches(0.12)) * i, gy + Inches(0.22), lwid2,
             Inches(0.42), lg, None, tcolor=WHITE,
             fill=RGBColor(0x14, 0x26, 0x50), tsize=9.5)
    tb, tf = textbox(s, Inches(10.05), gy + Inches(0.14), Inches(2.6), gh - Inches(0.28),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "Detects explicit requests and language cues, then swaps both voice "
             "and recognition instantly.", 9, MUTED, first=True, space_after=0)

    footer(s, "Stack: Python + FastAPI \u00b7 Twilio / Amazon Connect \u00b7 provider STT "
              "\u00b7 Sarvam AI / Amazon Polly \u00b7 AWS Bedrock (Claude) \u00b7 DynamoDB "
              "\u00b7 Amazon EC2")
    nav_bar(s, 4)
    add_transition(s)


# ================================================================== SLIDE 6
def slide6_value():
    s = prs.slides.add_slide(BLANK)
    fill_background(s)
    slide_head(s, "Combined Value",
               GREEN,
               "Two Platforms, One Engineering Standard",
               "What both systems share \u2014 and why that matters for your rollout",
               accent=GREEN)

    top = Inches(1.62)
    ch = Inches(2.10)
    gap = Inches(0.22)
    cw = (SW - Inches(1.1) - gap * 2) / 3
    x0 = Inches(0.55)

    card(s, x0, top, cw, ch, GREEN, "Built to scale", [
        "Thousands of calls in parallel, round the clock",
        "Agent layer adds specialists with no routing changes",
        "Data-driven catalogues and scripts \u2014 swap per client",
    ], center_heading=True)

    card(s, x0 + cw + gap, top, cw, ch, CYAN, "Secure by design", [
        "JWT + RBAC on every staff surface, 4 roles",
        "Encryption at rest (SSE-KMS), 1h presigned URLs",
        "Local LLM option \u2014 data never leaves the account",
    ], center_heading=True)

    card(s, x0 + (cw + gap) * 2, top, cw, ch, PURPLE, "Production-minded", [
        "Forward-only state machines, append-only audit trails",
        "Graceful degradation instead of 5xx failures",
        "Dockerized, CI to ECR, deployed on ECS / EKS",
    ], center_heading=True)

    # ---- shared stack strip ----
    sy = Inches(3.94)
    sh_ = Inches(0.94)
    rounded(s, x0, sy, SW - Inches(1.1), sh_, RGBColor(0x0E, 0x18, 0x32),
            line=BRD, line_w=1.0, radius=0.07)
    rect(s, x0 + Inches(0.02), sy + Inches(0.02), SW - Inches(1.14), Inches(0.055), BLUE)
    tb, tf = textbox(s, x0 + Inches(0.26), sy + Inches(0.13), Inches(2.4),
                     sh_ - Inches(0.26), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "Shared foundations", 12, BLUE, bold=True, first=True, space_after=0)
    chips = ["Python + FastAPI", "LLM orchestration", "PostgreSQL / DynamoDB",
             "AWS managed services", "Docker + CI/CD", "Prometheus / CloudWatch"]
    cwid = Inches(1.62)
    cxx = Inches(3.15)
    for i, chp in enumerate(chips):
        cell(s, cxx + (cwid + Inches(0.10)) * i, sy + Inches(0.26), cwid, Inches(0.42),
             chp, None, tcolor=INK, fill=RGBColor(0x14, 0x26, 0x50), tsize=8.8)

    # ---- closing band ----
    cy = Inches(5.06)
    cbh = Inches(1.26)
    rounded(s, x0, cy, SW - Inches(1.1), cbh, RGBColor(0x14, 0x14, 0x40),
            line=BRD_HI, line_w=1.25, radius=0.07)
    tb, tf = textbox(s, x0, cy + Inches(0.16), SW - Inches(1.1), cbh - Inches(0.32),
                     anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "Thank You", 30, WHITE, bold=True, align=PP_ALIGN.CENTER, first=True,
         space_after=4)
    para(tf, "Smarter onboarding and smarter conversations, at scale  \u00b7  "
             "Questions & Discussion", 13, CYAN, align=PP_ALIGN.CENTER, space_after=0)

    footer(s, "teleGlobal \u2014 Elevating business to technology  |  "
              "click any label in the bar below to revisit a section")
    nav_bar(s, 5)
    add_transition(s)


# ------------------------------------------------------------------ build
# Slides must all exist before nav links can resolve, so build the frames in
# order and attach navigation afterwards.
BUILDERS = [slide1_hub, slide2_bank_arch, slide3_bank_prod,
            slide4_voice_why, slide5_voice_how, slide6_value]


def build():
    # Pass 1: create the six slides. Clickable shapes only queue their target.
    for fn in BUILDERS:
        fn()
    # Pass 2: every slide now exists, so the slide-jump actions can be applied.
    return resolve_links()


if __name__ == "__main__":
    links = build()
    out = os.path.join(HERE, "teleGlobal_AI_Showcase_Interactive.pptx")
    prs.save(out)
    print("Logo:", LOGO or "not found (drew wordmark fallback)")
    print("Slides:", len(prs.slides._sldIdLst))
    print("Links:", links)
    print("Saved:", out)
